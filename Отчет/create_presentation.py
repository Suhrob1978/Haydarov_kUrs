# -*- coding: utf-8 -*-
"""Generate defense presentation for Kinopolka coursework."""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT = Path(__file__).resolve().parent.parent
ASSETS = Path(__file__).resolve().parent / "assets"
OUT = Path(__file__).resolve().parent / "Презентация_Кинополка.pptx"

NAVY = RGBColor(26, 54, 93)
BLUE = RGBColor(0, 102, 180)
GRAY = RGBColor(64, 64, 64)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(240, 244, 248)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W, H = prs.slide_width, prs.slide_height


def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_bar(slide, color=NAVY, height=Inches(1.05)):
    bar = slide.shapes.add_shape(1, 0, 0, W, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def title_text(slide, text, top=Inches(0.22), size=32, color=WHITE, width=None):
    box = slide.shapes.add_textbox(Inches(0.55), top, width or W - Inches(1.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = "Calibri"
    return box


def bullets(slide, items, left=Inches(0.7), top=Inches(1.35), width=Inches(12), size=22, spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, H - top - Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.name = "Calibri"
        p.font.color.rgb = GRAY
        p.space_after = Pt(size * spacing * 0.35)
        p.bullet = True
    return box


def section_slide(title, items, subtitle=None):
    slide = blank()
    add_bar(slide)
    title_text(slide, title, size=30)
    top = Inches(1.25)
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.7), top, W - Inches(1.4), Inches(0.5))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(18)
        sp.font.color.rgb = BLUE
        sp.font.name = "Calibri"
        top = Inches(1.75)
    bullets(slide, items, top=top)
    return slide


def image_slide(title, image_name, caption=None):
    slide = blank()
    add_bar(slide)
    title_text(slide, title, size=28)
    img_path = ASSETS / image_name
    if img_path.exists():
        slide.shapes.add_picture(str(img_path), Inches(0.55), Inches(1.15), width=Inches(12.2))
    else:
        ph = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), W - Inches(1.4), Inches(1))
        ph.text_frame.paragraphs[0].text = f"[{image_name} not found]"
    if caption:
        cap = slide.shapes.add_textbox(Inches(0.7), Inches(6.85), W - Inches(1.4), Inches(0.45))
        cp = cap.text_frame.paragraphs[0]
        cp.text = caption
        cp.font.size = Pt(14)
        cp.font.color.rgb = GRAY
        cp.font.name = "Calibri"
        cp.alignment = PP_ALIGN.CENTER
    return slide


def two_column_slide(title, left_items, right_items, left_title="", right_title=""):
    slide = blank()
    add_bar(slide)
    title_text(slide, title, size=28)
    if left_title:
        lt = slide.shapes.add_textbox(Inches(0.7), Inches(1.2), Inches(5.8), Inches(0.4))
        lp = lt.text_frame.paragraphs[0]
        lp.text = left_title
        lp.font.bold = True
        lp.font.size = Pt(20)
        lp.font.color.rgb = BLUE
        lp.font.name = "Calibri"
    if right_title:
        rt = slide.shapes.add_textbox(Inches(6.9), Inches(1.2), Inches(5.8), Inches(0.4))
        rp = rt.text_frame.paragraphs[0]
        rp.text = right_title
        rp.font.bold = True
        rp.font.size = Pt(20)
        rp.font.color.rgb = BLUE
        rp.font.name = "Calibri"
    bullets(slide, left_items, left=Inches(0.7), top=Inches(1.65), width=Inches(5.8), size=18)
    bullets(slide, right_items, left=Inches(6.9), top=Inches(1.65), width=Inches(5.8), size=18)
    return slide


# --- Slide 1: Title ---
s1 = blank()
bg = s1.shapes.add_shape(1, 0, 0, W, H)
bg.fill.solid()
bg.fill.fore_color.rgb = NAVY
bg.line.fill.background()
title_text(s1, "Курсовая работа", top=Inches(1.6), size=28, color=WHITE)
title_text(s1, "Разработка клиент-серверного\nмобильного приложения", top=Inches(2.15), size=34, color=WHITE)
title_text(s1, "«Бэклог фильмов» (Кинополка)", top=Inches(3.15), size=30, color=RGBColor(180, 210, 255))
info = s1.shapes.add_textbox(Inches(0.55), Inches(4.5), W - Inches(1.1), Inches(2.2))
tf = info.text_frame
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
lines = [
    "Дисциплина: Разработка кроссплатформенных и серверных мобильных приложений",
    "Выполнил: студент группы ИКБО-65-23  Хайдаров С.Р.",
    "Проверил: ____________________",
    "Москва, 2025",
]
for i, line in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = line
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.name = "Calibri"
    p.space_after = Pt(8)
    p.alignment = PP_ALIGN.CENTER

# --- Content slides ---
section_slide(
    "Цель и задачи работы",
    [
        "Цель: разработать и протестировать клиент-серверное Android-приложение для ведения личного бэклога фильмов.",
        "Проанализировать предметную область и существующие аналоги.",
        "Обосновать выбор технологического стека.",
        "Определить функциональные и нефункциональные требования.",
        "Спроектировать интерфейс, логическую модель и базу данных.",
        "Реализовать клиентскую и серверную части, провести тестирование.",
    ],
)

section_slide(
    "Актуальность и проблема",
    [
        "Списки «надо посмотреть» быстро разрастаются и теряются в заметках и мессенджерах.",
        "Стриминговые сервисы не дают удобного личного учёта с оценками и заметками.",
        "Пользователю нужен простой инструмент: статус просмотра, оценка, заметка, статистика.",
        "«Кинополка» решает задачу структурированного учёта фильмов с синхронизацией на сервере.",
        "Данные доступны с любого устройства после авторизации.",
    ],
)

section_slide(
    "Сравнение с аналогами",
    [
        "Letterboxd — социальная платформа, англоязычный интерфейс, нет открытого API.",
        "IMDb — огромная база, но личный учёт просмотра реализован ограниченно.",
        "Кинопоиск — ориентирован на продажу подписки, а не на личный бэклог.",
        "«Кинополка»: русскоязычный UI, фокус на статусах и заметках, открытая REST-архитектура.",
        "Клиент-серверное хранение данных и возможность дальнейшего развития.",
    ],
)

section_slide(
    "Функциональные возможности",
    [
        "Регистрация и вход с JWT-авторизацией.",
        "Каталог фильмов с поиском по названию (26 фильмов в демо-каталоге).",
        "Добавление в бэклог: «Хочу посмотреть», «Смотрю», «Посмотрел».",
        "Пользовательская оценка (0–10) и текстовые заметки.",
        "Фильтрация бэклога по статусу, статистика в профиле.",
        "Синхронизация данных между клиентом и сервером по REST API.",
    ],
)

two_column_slide(
    "Технологический стек",
    [
        "Kotlin",
        "Jetpack Compose, Material 3",
        "Navigation Compose, Hilt",
        "Retrofit, OkHttp, Coil",
        "ViewModel, StateFlow, MVVM",
        "Android SDK 35, minSdk 24",
    ],
    [
        "Kotlin, Spring Boot 3",
        "Spring Web, Spring Data JPA",
        "Hibernate, H2 (файловая БД)",
        "JJWT — авторизация",
        "BCrypt — хеширование паролей",
        "REST API, JSON",
    ],
    left_title="Android-клиент",
    right_title="Сервер",
)

image_slide(
    "Архитектура приложения",
    "arch.png",
    "Клиент (MVVM) ↔ REST/JSON/JWT ↔ Сервер (Controller → Service → Repository → H2)",
)

image_slide(
    "Диаграмма вариантов использования",
    "usecase.png",
    "Регистрация, работа с каталогом, управление бэклогом, просмотр профиля",
)

image_slide(
    "Логическая модель и БД",
    "er.png",
    "Таблицы: users, movies, backlog_items (связь многие-ко-многим через backlog_items)",
)

section_slide(
    "Пользовательский интерфейс",
    [
        "Три раздела нижней навигации: Каталог, Бэклог, Профиль.",
        "Каталог: поиск, карточки фильмов, добавление в бэклог.",
        "Бэклог: фильтры по статусу, редактирование оценки и заметки.",
        "Профиль: данные пользователя, статистика, выход из аккаунта.",
        "Material 3, поддержка светлой и тёмной темы.",
        "На слайдах отчёта — скриншоты работающего приложения.",
    ],
)

section_slide(
    "REST API (основные эндпоинты)",
    [
        "POST /api/auth/register — регистрация",
        "POST /api/auth/login — вход, получение JWT",
        "GET /api/movies?query= — каталог и поиск",
        "GET /api/backlog?status= — список бэклога",
        "POST /api/backlog — добавить фильм",
        "PATCH /api/backlog/{id} — изменить статус, оценку, заметку",
        "DELETE /api/backlog/{id} — удалить из бэклога",
        "GET /api/profile/stats — статистика пользователя",
    ],
    subtitle="Клиент обращается к http://10.0.2.2:8080 (эмулятор → localhost ПК)",
)

section_slide(
    "Запуск и демонстрация",
    [
        "1. Запустить сервер: start_server.bat или KinopolkaApplication в Android Studio.",
        "2. Дождаться строки Started KinopolkaApplicationKt (порт 8080).",
        "3. Открыть MovieBacklogApp в Android Studio, запустить на эмуляторе.",
        "4. Зарегистрироваться → найти фильм → добавить в бэклог.",
        "5. Изменить статус, поставить оценку, посмотреть статистику.",
        "Исходный код: github.com/Suhrob1978/HaydarovSR_KursovayaMobilki",
    ],
)

section_slide(
    "Тестирование",
    [
        "Функциональное (ручное) тестирование по сценариям из отчёта.",
        "Проверены: регистрация, вход, каталог, поиск, бэклог, профиль.",
        "Обработка ошибок: неверный пароль, повторное добавление, отсутствие сети.",
        "Сборка сервера с контекстом Spring Boot — успешно.",
        "Все запланированные сценарии пройдены.",
    ],
)

section_slide(
    "Заключение",
    [
        "Разработано клиент-серверное приложение «Кинополка» для Android.",
        "Реализованы учёт фильмов по статусам, оценки, заметки и статистика.",
        "Применены современные технологии: Compose, Spring Boot, JWT, H2.",
        "Приложение работоспособно и готово к демонстрации.",
        "Направления развития: рекомендации, интеграция с Kinopoisk/TMDb API, деплой сервера.",
    ],
)

# --- Final slide ---
s_end = blank()
add_bar(s_end, color=BLUE)
title_text(s_end, "Спасибо за внимание!", top=Inches(2.6), size=40, color=NAVY)
title_text(s_end, "Вопросы?", top=Inches(3.5), size=32, color=GRAY)
link = s_end.shapes.add_textbox(Inches(0.55), Inches(5.2), W - Inches(1.1), Inches(0.6))
lp = link.text_frame.paragraphs[0]
lp.text = "GitHub: github.com/Suhrob1978/HaydarovSR_KursovayaMobilki"
lp.font.size = Pt(18)
lp.font.color.rgb = BLUE
lp.font.name = "Calibri"
lp.alignment = PP_ALIGN.CENTER

prs.save(str(OUT))
print(f"Created: {OUT}")
print(f"Slides: {len(prs.slides)}")
